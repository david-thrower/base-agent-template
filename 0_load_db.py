import chromadb
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
import os
import glob
from typing import List, Optional
import re

TARGET_TOKENS = 150
PROJECT_NAME = "my_project"

DB_PATH = "/data/user/chroma_db"

# Document extraction libraries
try:
    import PyPDF2
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("Warning: PyPDF2 not installed. PDF support disabled. Install with: pip install PyPDF2")

try:
    import pdfplumber
    PDFPLUMBER_SUPPORT = True
except ImportError:
    PDFPLUMBER_SUPPORT = False

try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("Warning: python-docx not installed. Word support disabled. Install with: pip install python-docx")

try:
    import pandas as pd
    PANDAS_SUPPORT = True
except ImportError:
    PANDAS_SUPPORT = False
    print("Warning: pandas not installed. Excel/CSV support limited. Install with: pip install pandas")

try:
    import openpyxl  # noqa: F401
    OPENPYXL_SUPPORT = True
except ImportError:
    OPENPYXL_SUPPORT = False

try:
    from PIL import Image
    from pytesseract import image_to_string
    OCR_SUPPORT = True
except ImportError:
    OCR_SUPPORT = False
    print("Warning: pytesseract/PIL not installed. Image OCR disabled. Install with: pip install pytesseract pillow")

try:
    import pptx
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("Warning: python-pptx not installed. PowerPoint support disabled. Install with: pip install python-pptx")

# ==================== 1. SETUP VECTOR DB ====================
# Persistent storage in ./chroma_db folder
chroma_client = chromadb.PersistentClient(path=DB_PATH)

# Local embeddings (free, offline, fast for internal use)
# all-MiniLM-L6-v2 has a 256 token limit, but optimal chunking is typically 128-256 tokens
# Using ~200 tokens as a safe default for good semantic coherence
EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"
OPTIMAL_TOKENS_PER_CHUNK = 200  # Conservative default for this model
TOKENS_PER_WORD = 1.3  # Approximate ratio for English text

embedding_fn = SentenceTransformerEmbeddingFunction(
    model_name=EMBEDDING_MODEL_NAME  # 384-dim, good balance of speed/quality
)

collection = chroma_client.get_or_create_collection(
    name=f"{PROJECT_NAME}_knowledge",
    embedding_function=embedding_fn,
    metadata={"hnsw:space": "cosine"}
)


# ==================== 2. DOCUMENT TEXT EXTRACTION ====================

class DocumentExtractor:
    """Handles extraction of text from various document formats."""

    # Mapping of extensions to extraction methods
    SUPPORTED_EXTENSIONS = {
        '.txt': 'extract_text',
        '.md': 'extract_text',
        '.rst': 'extract_text',
        '.py': 'extract_text',
        '.js': 'extract_text',
        '.ts': 'extract_text',
        '.jsx': 'extract_text',
        '.tsx': 'extract_text',
        '.json': 'extract_text',
        '.yaml': 'extract_text',
        '.yml': 'extract_text',
        '.csv': 'extract_csv',
        '.log': 'extract_text',
        '.pdf': 'extract_pdf',
        '.docx': 'extract_docx',
        '.doc': 'extract_docx',  # Will try, may need antiword for .doc
        '.xlsx': 'extract_excel',
        '.xls': 'extract_excel',
        '.pptx': 'extract_pptx',
        '.ppt': 'extract_pptx',  # Will try
        '.png': 'extract_image',
        '.jpg': 'extract_image',
        '.jpeg': 'extract_image',
        '.tiff': 'extract_image',
        '.bmp': 'extract_image',
        '.gif': 'extract_image',
    }

    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extract text from plain text files."""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Could not decode file: {file_path}")

    @staticmethod
    def extract_pdf(file_path: str) -> str:  # noqa: C901
        """Extract text from PDF using pdfplumber (preferred) or PyPDF2."""
        text = ""

        # Try pdfplumber first (better for complex layouts)
        if PDFPLUMBER_SUPPORT:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                if text.strip():
                    return text
            except Exception as e:
                print(f"pdfplumber failed for {file_path}: {e}")

        # Fallback to PyPDF2
        if PDF_SUPPORT:
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() or ""
                return text
            except Exception as e:
                raise ValueError(f"Failed to extract PDF {file_path}: {e}")
        else:
            raise ValueError("No PDF extraction library available")

    @staticmethod
    def extract_docx(file_path: str) -> str:
        """Extract text from Word documents."""
        if not DOCX_SUPPORT:
            raise ValueError("python-docx not installed")

        try:
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

            # Also extract text from tables
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells if cell.text.strip()]
                    if row_text:
                        tables_text.append(" | ".join(row_text))

            all_text = "\n".join(paragraphs)
            if tables_text:
                all_text += "\n\n[Tables]\n" + "\n".join(tables_text)

            return all_text
        except Exception as e:
            raise ValueError(f"Failed to extract DOCX {file_path}: {e}")

    @staticmethod
    def extract_excel(file_path: str) -> str:
        """Extract text from Excel files as structured text."""
        if not PANDAS_SUPPORT:
            raise ValueError("pandas not installed")

        try:
            # Read all sheets
            xl_file = pd.ExcelFile(file_path)
            all_sheets_text = []

            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)

                # Convert dataframe to string representation
                sheet_text = f"[Sheet: {sheet_name}]\n"
                sheet_text += df.to_string(index=False)
                all_sheets_text.append(sheet_text)

            return "\n\n".join(all_sheets_text)
        except Exception as e:
            raise ValueError(f"Failed to extract Excel {file_path}: {e}")

    @staticmethod
    def extract_csv(file_path: str) -> str:
        """Extract text from CSV files."""
        if not PANDAS_SUPPORT:
            # Fallback to basic text reading
            return DocumentExtractor.extract_text(file_path)

        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception:
            # Fallback to basic text reading if pandas fails
            return DocumentExtractor.extract_text(file_path)

    @staticmethod
    def extract_pptx(file_path: str) -> str:
        """Extract text from PowerPoint presentations."""
        if not PPTX_SUPPORT:
            raise ValueError("python-pptx not installed")

        try:
            prs = pptx.Presentation(file_path)
            all_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if len(slide_text) > 1:
                    all_text.append("\n".join(slide_text))

            return "\n\n".join(all_text)
        except Exception as e:
            raise ValueError(f"Failed to extract PPTX {file_path}: {e}")

    @staticmethod
    def extract_image(file_path: str) -> str:
        """Extract text from images using OCR."""
        if not OCR_SUPPORT:
            raise ValueError("OCR libraries not installed (pytesseract, PIL)")

        try:
            image = Image.open(file_path)

            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')

            # Perform OCR
            text = image_to_string(image)
            return text
        except Exception as e:
            raise ValueError(f"Failed to OCR image {file_path}: {e}")

    @classmethod
    def extract(cls, file_path: str) -> str:
        """
        Extract text from any supported file type.

        Args:
            file_path: Path to the file

        Returns:
            Extracted text as string

        Raises:
            ValueError: If file type not supported or extraction fails
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext not in cls.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {ext}")

        method_name = cls.SUPPORTED_EXTENSIONS[ext]
        method = getattr(cls, method_name)

        return method(file_path)


# ==================== 3. DOCUMENT INGESTION PIPELINE ====================

def estimate_tokens(text: str) -> int:
    """Rough estimation of token count based on word count."""
    word_count = len(text.split())
    return int(word_count * TOKENS_PER_WORD)


def chunk_text(text: str, target_tokens: int = OPTIMAL_TOKENS_PER_CHUNK) -> List[str]:
    """
    Split text into chunks of approximately target_tokens.
    Uses sentence-aware splitting to preserve semantic coherence.
    """
    if not text or not text.strip():
        return []

    words = text.split()
    target_words = int(target_tokens / TOKENS_PER_WORD)

    chunks = []
    current_chunk = []
    current_word_count = 0

    for word in words:
        current_chunk.append(word)
        current_word_count += 1

        # Check if we've reached target size and are at a sentence boundary
        if current_word_count >= target_words and word.endswith(('.', '!', '?')):
            chunks.append(" ".join(current_chunk))
            current_chunk = []
            current_word_count = 0

    # Add any remaining words
    if current_chunk:
        chunks.append(" ".join(current_chunk))

    # If no chunks were created (e.g., no sentence boundaries), just split by word count
    if not chunks and words:
        for i in range(0, len(words), target_words):
            chunk_words = words[i:i + target_words]
            chunks.append(" ".join(chunk_words))

    return chunks


def ingest_documents(  # noqa: C901
    documents_dir: str = "./documents",
    target_tokens: int = OPTIMAL_TOKENS_PER_CHUNK,
    file_extensions: Optional[List[str]] = None
) -> dict:
    """
    Traverse documents directory, extract text from all files, chunk,
    and add to vector DB.

    Args:
        documents_dir: Directory containing documents
        target_tokens: Target tokens per chunk
        file_extensions: Optional list of specific extensions to process

    Returns:
        dict: Statistics about ingestion process
    """
    if not os.path.exists(documents_dir):
        print(f"Documents directory '{documents_dir}' does not exist. Skipping ingestion.")
        return {"status": "skipped", "reason": "directory_not_found"}

    # Determine which extensions to process
    if file_extensions:
        extensions_to_process = file_extensions
    else:
        extensions_to_process = list(DocumentExtractor.SUPPORTED_EXTENSIONS.keys())

    stats = {
        "files_processed": 0,
        "files_skipped": 0,
        "chunks_created": 0,
        "total_documents_before": collection.count(),
        "errors": [],
        "by_type": {}
    }

    # Find all files recursively
    all_files = []
    for ext in extensions_to_process:
        pattern = os.path.join(documents_dir, f"**/*{ext}")
        all_files.extend(glob.glob(pattern, recursive=True))

    # Remove duplicates and sort
    all_files = sorted(set(all_files))

    print(f"Found {len(all_files)} files to process in '{documents_dir}'")
    print(f"Supported types: {', '.join(extensions_to_process)}")

    for file_path in all_files:
        try:
            # Get relative path and file info
            rel_path = os.path.relpath(file_path, documents_dir)
            file_ext = os.path.splitext(file_path)[1].lower()

            print(f"Processing: {rel_path}")

            # Extract text using appropriate method
            content = DocumentExtractor.extract(file_path)

            if not content or not content.strip():
                print(f"  ⚠ No content extracted from {rel_path}")
                stats["files_skipped"] += 1
                continue

            # Clean up whitespace
            content = re.sub(r'\s+', ' ', content).strip()

            # Split into chunks
            chunks = chunk_text(content, target_tokens)

            if not chunks:
                print(f"  ⚠ No chunks created from {rel_path}")
                stats["files_skipped"] += 1
                continue

            # Prepare metadata for each chunk
            metadatas = []
            ids = []

            for i, chunk in enumerate(chunks):
                metadata = {
                    "source": rel_path,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "file_type": file_ext,
                    "estimated_tokens": estimate_tokens(chunk),
                    "extraction_method": DocumentExtractor.SUPPORTED_EXTENSIONS.get(file_ext, "unknown")
                }
                metadatas.append(metadata)
                # Create unique ID: filename_chunk_index
                safe_filename = re.sub(r'[^a-zA-Z0-9_-]', '_', os.path.splitext(rel_path)[0])
                chunk_id = f"{safe_filename}_{i}_{int(os.path.getmtime(file_path))}"
                ids.append(chunk_id)

            # Add to collection
            collection.add(
                documents=chunks,
                ids=ids,
                metadatas=metadatas
            )

            stats["files_processed"] += 1
            stats["chunks_created"] += len(chunks)

            # Track by file type
            if file_ext not in stats["by_type"]:
                stats["by_type"][file_ext] = {"count": 0, "chunks": 0}
            stats["by_type"][file_ext]["count"] += 1
            stats["by_type"][file_ext]["chunks"] += len(chunks)

            print(f"  ✓ {len(chunks)} chunks created (~{estimate_tokens(chunks[0])} tokens each)")

        except Exception as e:
            error_msg = f"Error processing {file_path}: {str(e)}"
            stats["errors"].append(error_msg)
            print(f"  ✗ {error_msg}")

    stats["total_documents_after"] = collection.count()
    stats["documents_added"] = stats["total_documents_after"] - stats["total_documents_before"]

    # Print summary
    print(f"\n{'='*50}")
    print(f"INGESTION COMPLETE")
    print(f"{'='*50}")
    print(f"Files processed: {stats['files_processed']}")
    print(f"Files skipped: {stats['files_skipped']}")
    print(f"Chunks created: {stats['chunks_created']}")
    print(f"Total documents in DB: {stats['total_documents_after']}")

    if stats["by_type"]:
        print(f"\nBreakdown by type:")
        for ext, data in sorted(stats["by_type"].items()):
            print(f"  {ext}: {data['count']} files, {data['chunks']} chunks")

    if stats["errors"]:
        print(f"\nErrors ({len(stats['errors'])}):")
        for err in stats["errors"][:5]:  # Show first 5 errors
            print(f"  - {err}")

    return stats


def clear_collection():
    """Clear all documents from the collection."""
    try:
        chroma_client.delete_collection("project_knowledge")
        global collection
        collection = chroma_client.get_or_create_collection(
            name="project_knowledge",
            embedding_function=embedding_fn,
            metadata={"hnsw:space": "cosine"}
        )
        print("Collection cleared successfully.")
    except Exception as e:
        print(f"Error clearing collection: {e}")


def search_documents(query: str, n_results: int = 5) -> List[dict]:
    """Search the vector database for relevant documents."""
    results = collection.query(
        query_texts=[query],
        n_results=n_results,
        include=["documents", "metadatas", "distances"]
    )

    documents = []
    for i in range(len(results['ids'][0])):
        documents.append({
            'id': results['ids'][0][i],
            'text': results['documents'][0][i],
            'metadata': results['metadatas'][0][i],
            'distance': results['distances'][0][i]
        })

    return documents


# Seed with sample data (run once, then comment out)
# def init_knowledge():
#     docs = [
#         "Project Phoenix uses React frontend with FastAPI backend",
#         "Authentication is handled via SSO with Azure AD",
#         "Database schema: Users table has id, email, role, created_at",
#         "Deployment pipeline: GitHub Actions → Docker → AWS ECS",
#         "API rate limits: 1000 requests/hour per API key"
#     ]
#     collection.add(
#         documents=docs,
#         ids=[f"doc_{i}" for i in range(len(docs))],
#         metadatas=[{"source": "internal_wiki"} for _ in docs]
#     )


# ==================== 4. RUN INGESTION ====================
if __name__ == "__main__":
    # Uncomment to run ingestion on startup:

    # Or with custom settings:
    ingest_documents(target_tokens=TARGET_TOKENS)
    # ingest_documents(file_extensions=['.pdf', '.docx'])  # Only specific types
